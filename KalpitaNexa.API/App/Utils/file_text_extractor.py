# app/Utils/file_text_extractor.py
"""
This utility module provides functions to extract plain text content from
various file formats like PDF, DOCX, PPTX, XLSX, TXT, and images (PNG).
"""
import logging
import io
from fastapi import UploadFile

# Import necessary libraries for each file type
try:
    import docx
except ImportError:
    docx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None
    
try:
    from PIL import Image
    import pytesseract

except ImportError:
    Image = None
    pytesseract = None

logger = logging.getLogger(__name__)

def extract_text_from_file(file: UploadFile) -> str:
    """
    Dispatcher function that reads an uploaded file's content and routes it
    to the correct text extraction helper based on its file extension.

    Args:
        file: The UploadFile object from FastAPI.

    Returns:
        The extracted plain text as a string, or an error message if extraction fails.
    """
    file_extension = f".{file.filename.split('.')[-1].lower()}"
    
    try:
        content = file.file.read()
    finally:
        file.file.seek(0)

    logger.info(f"Extracting text from '{file.filename}' (type: {file_extension})")
    
    try:
        if file_extension == '.docx':
            return _extract_text_from_docx(content)
        elif file_extension == '.pdf':
            return _extract_text_from_pdf(content)
        elif file_extension in ['.ppt', '.pptx']:
            return _extract_text_from_pptx(content)
        elif file_extension == '.xlsx':
            return _extract_text_from_xlsx(content)
        elif file_extension == '.txt':
            return _extract_text_from_txt(content)
        elif file_extension == '.png':
            return _extract_text_from_png(content)
        else:
            return f"Error: Unsupported file type '{file_extension}'."
    except Exception as e:
        logger.error(f"Failed to extract text from {file.filename}: {e}", exc_info=True)
        return f"Error: Could not process file {file.filename}."

def _extract_text_from_docx(content: bytes) -> str:
    """Extracts text from a .docx file."""
    if not docx:
        raise ImportError("The 'python-docx' library is not installed.")
    doc = docx.Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs if para.text])

def _extract_text_from_pdf(content: bytes) -> str:
    """Extracts text from a .pdf file."""
    if not PyPDF2:
        raise ImportError("The 'PyPDF2' library is not installed.")
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    return "\n".join([page.extract_text() for page in reader.pages])

def _extract_text_from_pptx(content: bytes) -> str:
    """Extracts text from a .pptx file."""
    if not pptx:
        raise ImportError("The 'python-pptx' library is not installed.")
    presentation = pptx.Presentation(io.BytesIO(content))
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def _extract_text_from_xlsx(content: bytes) -> str:
    """Extracts text from an .xlsx file, cell by cell."""
    if not openpyxl:
        raise ImportError("The 'openpyxl' library is not installed.")
    workbook = openpyxl.load_workbook(io.BytesIO(content))
    text_content = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            row_text = [str(cell.value) for cell in row if cell.value is not None]
            if row_text:
                text_content.append(" | ".join(row_text))
    return "\n".join(text_content)

def _extract_text_from_txt(content: bytes) -> str:
    """Extracts text from a .txt file."""
    try:
        # Try decoding with utf-8, but fall back to latin-1 for broader compatibility
        return content.decode('utf-8')
    except UnicodeDecodeError:
        return content.decode('latin-1', errors='ignore')

def _extract_text_from_png(content: bytes) -> str:
    """Extracts text from a .png image file using OCR (Tesseract)."""
    if not Image or not pytesseract:
        raise ImportError("The 'Pillow' and 'pytesseract' libraries are required for image processing.")
    
    try:
        image = Image.open(io.BytesIO(content))
        # Use pytesseract to do OCR on the image
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        logger.error(f"OCR processing failed: {e}")
        return "Error: OCR processing failed. Ensure Tesseract is installed and configured correctly."